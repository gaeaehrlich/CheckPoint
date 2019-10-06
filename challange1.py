from pptx import Presentation
import os.path

prs = Presentation('START.pptx')
sp = prs.slides[0].shapes[0]
params = sp.text.split(', ')
res = params[0]
while os.path.exists(params[1]):
    prs = Presentation(params[1])
    sp = prs.slides[int(params[2])-1].shapes[0]
    params = sp.text.split(', ')
    res += params[0]
print(res)
